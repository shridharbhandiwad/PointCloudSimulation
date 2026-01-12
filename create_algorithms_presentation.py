"""
Create a PowerPoint presentation focused on Radar Detection Algorithms
- Radar Detection Simulation
- DBSCAN Clustering
- Object Tracking
- Classification (Rule-Based & Naive Bayes)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(25, 55, 95)
ACCENT_BLUE = RGBColor(41, 128, 185)
LIGHT_BLUE = RGBColor(174, 214, 241)
ORANGE = RGBColor(230, 126, 34)
GREEN = RGBColor(39, 174, 96)
RED = RGBColor(192, 57, 43)
PURPLE = RGBColor(142, 68, 173)
GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_frame = subtitle_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(28)
        sub_para.font.color.rgb = LIGHT_BLUE
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, section_title, section_number=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = ACCENT_BLUE
    background.line.fill.background()
    
    # Section number
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = section_number
        num_para.font.size = Pt(36)
        num_para.font.color.rgb = LIGHT_BLUE
        num_para.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = section_title
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items, two_column=False):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    if two_column and len(content_items) == 2:
        # Two columns
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        col1_title, col1_items = content_items[0]
        p = left_frame.paragraphs[0]
        p.text = col1_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        for item in col1_items:
            p = left_frame.add_paragraph()
            if isinstance(item, tuple):
                p.text = f"• {item[0]}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = BLACK
                p.space_before = Pt(8)
                for sub in item[1]:
                    sp = left_frame.add_paragraph()
                    sp.text = f"   – {sub}"
                    sp.font.size = Pt(16)
                    sp.font.color.rgb = GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = BLACK
                p.space_before = Pt(8)
        
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        col2_title, col2_items = content_items[1]
        p = right_frame.paragraphs[0]
        p.text = col2_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        for item in col2_items:
            p = right_frame.add_paragraph()
            if isinstance(item, tuple):
                p.text = f"• {item[0]}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = BLACK
                p.space_before = Pt(8)
                for sub in item[1]:
                    sp = right_frame.add_paragraph()
                    sp.text = f"   – {sub}"
                    sp.font.size = Pt(16)
                    sp.font.color.rgb = GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = BLACK
                p.space_before = Pt(8)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        first = True
        for item in content_items:
            if first:
                p = content_frame.paragraphs[0]
                first = False
            else:
                p = content_frame.add_paragraph()
            
            if isinstance(item, tuple):
                p.text = f"• {item[0]}"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = BLACK
                p.space_before = Pt(12)
                for sub in item[1]:
                    sp = content_frame.add_paragraph()
                    sp.text = f"   – {sub}"
                    sp.font.size = Pt(18)
                    sp.font.color.rgb = GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(20)
                p.font.color.rgb = BLACK
                p.space_before = Pt(12)
    
    return slide


def add_algorithm_slide(prs, title, algorithm_name, formula, description, steps):
    """Add a slide showcasing an algorithm with formula and steps."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Algorithm name box
    algo_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4), Inches(0.6)
    )
    algo_box.fill.solid()
    algo_box.fill.fore_color.rgb = ACCENT_BLUE
    algo_box.line.fill.background()
    
    algo_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(4), Inches(0.5))
    atf = algo_text.text_frame
    ap = atf.paragraphs[0]
    ap.text = algorithm_name
    ap.font.size = Pt(22)
    ap.font.bold = True
    ap.font.color.rgb = WHITE
    ap.alignment = PP_ALIGN.CENTER
    
    # Formula box
    if formula:
        formula_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.4), Inches(7.8), Inches(0.8)
        )
        formula_shape.fill.solid()
        formula_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        formula_shape.line.color.rgb = GRAY
        
        formula_text = slide.shapes.add_textbox(Inches(5.1), Inches(1.5), Inches(7.6), Inches(0.6))
        ftf = formula_text.text_frame
        fp = ftf.paragraphs[0]
        fp.text = formula
        fp.font.size = Pt(20)
        fp.font.name = "Consolas"
        fp.font.color.rgb = DARK_BLUE
        fp.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.333), Inches(0.8))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = description
    dp.font.size = Pt(18)
    dp.font.color.rgb = GRAY
    
    # Steps
    steps_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(4))
    stf = steps_box.text_frame
    stf.word_wrap = True
    
    first = True
    for i, step in enumerate(steps, 1):
        if first:
            sp = stf.paragraphs[0]
            first = False
        else:
            sp = stf.add_paragraph()
        
        sp.text = f"{i}. {step}"
        sp.font.size = Pt(18)
        sp.font.color.rgb = BLACK
        sp.space_before = Pt(10)
    
    return slide


def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a comparison slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = GREEN
    left_header.line.fill.background()
    
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(6), Inches(0.5))
    ltf = left_title_box.text_frame
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(22)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.alignment = PP_ALIGN.CENTER
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(6), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    first = True
    for item in left_items:
        if first:
            p = left_frame.paragraphs[0]
            first = False
        else:
            p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
    
    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.833), Inches(1.5), Inches(6), Inches(0.6)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = PURPLE
    right_header.line.fill.background()
    
    right_title_box = slide.shapes.add_textbox(Inches(6.833), Inches(1.55), Inches(6), Inches(0.5))
    rtf = right_title_box.text_frame
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(22)
    rp.font.bold = True
    rp.font.color.rgb = WHITE
    rp.alignment = PP_ALIGN.CENTER
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.833), Inches(2.3), Inches(6), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    first = True
    for item in right_items:
        if first:
            p = right_frame.paragraphs[0]
            first = False
        else:
            p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
    
    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a data table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table_width = Inches(12.333)
    table_height = Inches(0.5) * num_rows
    left = Inches(0.5)
    top = Inches(1.6)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
    
    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = BLACK
            para.alignment = PP_ALIGN.CENTER
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide


def add_flowchart_slide(prs, title, steps):
    """Add a slide with a horizontal flowchart."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Calculate positions
    num_steps = len(steps)
    box_width = Inches(2.2)
    box_height = Inches(1.2)
    arrow_width = Inches(0.6)
    total_width = num_steps * box_width + (num_steps - 1) * arrow_width
    start_x = (prs.slide_width - total_width) / 2
    y_pos = Inches(2.5)
    
    colors = [ACCENT_BLUE, GREEN, ORANGE, PURPLE, RED]
    
    for i, (step_title, step_desc) in enumerate(steps):
        x_pos = start_x + i * (box_width + arrow_width)
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        
        # Step number
        num_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.1), box_width, Inches(0.4))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = f"Step {i + 1}"
        np.font.size = Pt(14)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER
        
        # Step title
        title_tb = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.45), box_width, Inches(0.7))
        ttf = title_tb.text_frame
        tp = ttf.paragraphs[0]
        tp.text = step_title
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        tp.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < num_steps - 1:
            arrow_x = x_pos + box_width + Inches(0.05)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + Inches(0.4), Inches(0.5), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()
        
        # Description below
        desc_box = slide.shapes.add_textbox(x_pos, y_pos + box_height + Inches(0.2), box_width, Inches(2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = step_desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = GRAY
        dp.alignment = PP_ALIGN.CENTER
    
    return slide


# =============================================================================
# CREATE THE PRESENTATION
# =============================================================================

# Slide 1: Title
add_title_slide(
    prs,
    "Radar Detection Algorithms",
    "Simulation, Clustering, Tracking & Classification"
)

# Slide 2: Agenda
add_content_slide(prs, "Agenda", [
    ("Radar Detection Simulation", [
        "Radar equation and received power modeling",
        "Point cloud generation algorithms",
        "Noise modeling and measurement simulation"
    ]),
    ("DBSCAN Clustering Algorithm", [
        "Density-based spatial clustering",
        "Core points, border points, and noise",
        "Implementation and parameter tuning"
    ]),
    ("Multi-Object Tracking", [
        "Nearest-neighbor data association",
        "Track lifecycle management",
        "Kalman filter prediction (linear motion)"
    ]),
    ("Classification Algorithms", [
        "Rule-Based classifier with feature thresholds",
        "Naive Bayes probabilistic classification",
        "Model comparison and evaluation"
    ])
])

# Section 1: Radar Detection Simulation
add_section_slide(prs, "Radar Detection Simulation", "SECTION 1")

# Slide 4: Radar Equation
add_algorithm_slide(
    prs,
    "Radar Equation - Received Power",
    "Radar Equation",
    "Pr(dBm) = K + 10·log₁₀(σ) - 40·log₁₀(R)",
    "The radar equation calculates received signal power based on target RCS and range. K is a calibration constant (default: 81.2 dB).",
    [
        "Calculate range R from radar position to target (Euclidean distance in 3D)",
        "Get target RCS (σ) based on object type and surface properties",
        "Apply radar equation with path loss (-40·log₁₀(R) accounts for R⁴ loss)",
        "Add Gaussian noise (shadow fading) with σ = 1.2 dB for realistic variation",
        "Combine individual detection powers using linear summation (mW → dBm)"
    ]
)

# Slide 5: Point Cloud Generation
add_content_slide(prs, "Point Cloud Generation Algorithm", [
    ("Surface Sampling Strategy", [
        "Sample points on visible surfaces (front/back, sides, top)",
        "Weight sampling towards radar-facing surfaces",
        "Apply corner reflector effect (2x RCS at corners)"
    ]),
    ("Distance-Based Scaling", [
        "Closer objects → more detections (inverse relationship)",
        "At 10m: full detection count",
        "At 100m: ~25% of maximum count"
    ]),
    ("Per-Point RCS Distribution", [
        "Total object RCS distributed across detection points",
        "Corner points receive higher RCS values",
        "Random variation factor (0.5x - 1.5x) for realism"
    ]),
    ("Noise Application", [
        "Range noise: σ_r = 0.25m (Gaussian)",
        "Azimuth noise: σ_az = 0.2° (Gaussian)",
        "Radial velocity noise: σ_vr = 0.15 m/s (Gaussian)"
    ])
])

# Slide 6: Detection Count Table
add_table_slide(
    prs,
    "Detection Count by Object Type",
    ["Object Type", "Length (m)", "Width (m)", "RCS σ (m²)", "Min Detections", "Max Detections"],
    [
        ["Car", "4.5", "1.9", "10.0", "15", "40"],
        ["Truck", "10.0", "2.6", "30.0", "30", "80"],
        ["Two-Wheeler", "2.2", "0.8", "3.0", "5", "12"],
        ["Bicycle", "1.8", "0.6", "1.0", "2", "6"],
        ["Pedestrian", "0.5", "0.5", "0.5", "1", "4"],
    ]
)

# Section 2: DBSCAN Clustering
add_section_slide(prs, "DBSCAN Clustering Algorithm", "SECTION 2")

# Slide 8: DBSCAN Overview
add_algorithm_slide(
    prs,
    "DBSCAN - Density-Based Spatial Clustering",
    "DBSCAN Algorithm",
    "ε-neighborhood: N_ε(p) = {q ∈ D | dist(p,q) ≤ ε}",
    "DBSCAN groups points based on density, separating clusters of arbitrary shape from noise without requiring the number of clusters a priori.",
    [
        "For each unvisited point p, retrieve its ε-neighborhood",
        "If |N_ε(p)| ≥ minPts, p is a core point - start new cluster",
        "Recursively expand cluster by adding all density-reachable points",
        "If |N_ε(p)| < minPts and not reachable from core, mark as noise",
        "Continue until all points are visited and labeled"
    ]
)

# Slide 9: DBSCAN Point Classification
add_content_slide(prs, "DBSCAN Point Classification", [
    ("Core Points", [
        "Has at least minPts neighbors within ε distance",
        "Forms the dense interior of clusters",
        "Can expand the cluster by connecting to other core points"
    ]),
    ("Border Points", [
        "Within ε of a core point but < minPts neighbors",
        "Belongs to cluster but cannot expand it",
        "Assigned to nearest core point's cluster"
    ]),
    ("Noise Points", [
        "Not within ε of any core point",
        "Isolated detections (clutter, false alarms)",
        "Excluded from object tracking and classification"
    ]),
    ("Parameters", [
        "ε (epsilon): 2.0m default - maximum neighbor distance",
        "minPts: 2 default - minimum samples for core point"
    ])
])

# Slide 10: DBSCAN Implementation
add_content_slide(prs, "DBSCAN Implementation Details", [
    ("Distance Computation", [
        "Euclidean distance: d(p,q) = √[(x₁-x₂)² + (y₁-y₂)²]",
        "Applied to world coordinates (x_w, y_w)",
        "O(n²) naive implementation, O(n log n) with spatial indexing"
    ]),
    ("Cluster Expansion (BFS)", [
        "Use queue to track points to process",
        "When core point found, add all neighbors to queue",
        "Mark visited points to avoid reprocessing",
        "Continue until queue empty"
    ]),
    ("Feature Extraction per Cluster", [
        "Centroid: mean(x), mean(y) of all points",
        "Extent: (max-min) in X and Y directions",
        "Statistics: mean/std of RCS and radial velocity",
        "Point count: number of detections in cluster"
    ])
])

# Section 3: Multi-Object Tracking
add_section_slide(prs, "Multi-Object Tracking Algorithm", "SECTION 3")

# Slide 12: Tracking Overview
add_flowchart_slide(prs, "Tracking Pipeline", [
    ("Prediction", "Apply motion model to predict next state for all tracks"),
    ("Association", "Match clusters to tracks using distance metric"),
    ("Update", "Update matched tracks with cluster measurements"),
    ("Management", "Create new tracks, delete old tracks")
])

# Slide 13: Data Association Algorithm
add_algorithm_slide(
    prs,
    "Nearest-Neighbor Data Association",
    "Greedy Association",
    "cost(c,t) = ||centroid(c) - pos(t)||₂",
    "Association matches new cluster measurements to existing tracks using distance-based cost minimization.",
    [
        "Compute cost matrix: distance from each cluster to each track",
        "Sort all (cluster, track) pairs by cost (ascending)",
        "For each pair in sorted order:",
        "   - If cluster and track are both unassigned AND cost ≤ threshold",
        "   - Assign cluster to track, mark both as used",
        "Unassigned clusters → create new tracks",
        "Unassigned tracks → increment miss counter"
    ]
)

# Slide 14: Motion Model
add_content_slide(prs, "Motion Model - Linear Prediction", [
    ("State Vector", [
        "Position: (x, y) in world coordinates",
        "Velocity: (vx, vy) estimated from position changes"
    ]),
    ("Prediction Step", [
        "x(t+dt) = x(t) + vx · dt",
        "y(t+dt) = y(t) + vy · dt",
        "Constant velocity assumption (no acceleration)"
    ]),
    ("Velocity Estimation", [
        "vx = (x_new - x_old) / dt",
        "vy = (y_new - y_old) / dt",
        "Updated on each successful cluster association"
    ]),
    ("Track Lifecycle Parameters", [
        "min_hits = 2: confirmations needed for valid track",
        "max_misses = 5: consecutive misses before deletion",
        "association_threshold = 5.0m: max distance for matching"
    ])
])

# Slide 15: Track Lifecycle
add_content_slide(prs, "Track Lifecycle Management", [
    ("Track Creation", [
        "New track created for each unassigned cluster",
        "Initial state: cluster centroid, zero velocity",
        "Status: tentative (hits = 1, age = 1)"
    ]),
    ("Track Confirmation", [
        "Track confirmed when hits ≥ min_hits",
        "Only confirmed tracks used for classification",
        "Provides temporal consistency"
    ]),
    ("Track Update", [
        "Position updated from associated cluster centroid",
        "Velocity estimated from position delta",
        "Feature history maintained (last 20 frames)"
    ]),
    ("Track Deletion", [
        "Miss counter incremented on no association",
        "Track deleted when misses > max_misses",
        "Prevents ghost tracks from persisting"
    ])
])

# Section 4: Classification Algorithms
add_section_slide(prs, "Classification Algorithms", "SECTION 4")

# Slide 17: Classification Features
add_table_slide(
    prs,
    "Classification Features",
    ["Feature", "Description", "Discriminative Power"],
    [
        ["num_points", "Number of detections in cluster", "High - varies by object size"],
        ["extent_x", "Width in X direction (m)", "High - vehicle length indicator"],
        ["extent_y", "Width in Y direction (m)", "Medium - vehicle width"],
        ["mean_rcs", "Average RCS in dBm", "High - object reflectivity"],
        ["std_rcs", "RCS standard deviation", "Low - surface uniformity"],
        ["speed", "Estimated velocity (m/s)", "High - motion characteristics"],
        ["mean_vr", "Mean radial velocity", "Medium - approach/departure"],
    ]
)

# Slide 18: Rule-Based Classifier
add_algorithm_slide(
    prs,
    "Rule-Based Classification",
    "Threshold Matching",
    "score(c) = mean(match_scores) where match ∈ [0,1]",
    "Classification based on predefined feature ranges for each class. Fast, interpretable, requires domain knowledge.",
    [
        "For each class, define (min, max) range for each feature",
        "For input features, compute match score per feature:",
        "   - score = 1.0 if value in range (penalty for distance from center)",
        "   - score = max(0, 1 - distance/range_size) if outside",
        "Average scores across features for class score",
        "Select class with highest score",
        "Confidence = best_score / sum(all_scores)"
    ]
)

# Slide 19: Rule-Based Thresholds
add_table_slide(
    prs,
    "Rule-Based Classification Thresholds",
    ["Class", "Points", "Extent X (m)", "Extent Y (m)", "RCS (dBm)", "Speed (m/s)"],
    [
        ["Car", "10-50", "1.5-6.0", "1.0-3.0", "10-35", "0.5-40"],
        ["Truck", "25-200", "3.0-15.0", "1.5-5.0", "20-50", "0.5-30"],
        ["Two-Wheeler", "3-15", "0.8-3.0", "0.3-1.5", "5-25", "0.5-35"],
        ["Bicycle", "1-8", "0.5-2.5", "0.2-1.2", "-5-20", "0.3-15"],
        ["Pedestrian", "1-6", "0.2-1.5", "0.2-1.0", "-10-15", "0.1-3"],
        ["Clutter", "1-5", "0.0-1.0", "0.0-1.0", "-20-20", "0.0-0.5"],
    ]
)

# Slide 20: Naive Bayes Classifier
add_algorithm_slide(
    prs,
    "Naive Bayes Classification",
    "Probabilistic Model",
    "P(C|X) ∝ P(C) · ∏ P(xᵢ|C)",
    "Probabilistic classifier assuming feature independence. Outputs confidence probabilities for each class.",
    [
        "Assume features are conditionally independent given class",
        "Model each feature with Gaussian distribution: P(x|C) = N(μ_C, σ_C)",
        "Compute log-likelihood for each class:",
        "   log P(C|X) = log P(C) + Σ log P(xᵢ|C)",
        "Apply log-sum-exp trick for numerical stability",
        "Normalize to get posterior probabilities",
        "Select class with maximum posterior probability"
    ]
)

# Slide 21: Naive Bayes Parameters
add_table_slide(
    prs,
    "Naive Bayes Gaussian Parameters (μ, σ)",
    ["Class", "Points", "Extent X", "RCS (dBm)", "Speed (m/s)", "Prior P(C)"],
    [
        ["Car", "(25, 10)", "(3.5, 1.0)", "(22, 8)", "(12, 6)", "0.35"],
        ["Truck", "(50, 20)", "(8.0, 2.5)", "(30, 8)", "(10, 5)", "0.15"],
        ["Two-Wheeler", "(8, 4)", "(2.0, 0.6)", "(15, 7)", "(10, 5)", "0.15"],
        ["Bicycle", "(4, 2)", "(1.5, 0.5)", "(8, 6)", "(5, 3)", "0.10"],
        ["Pedestrian", "(2.5, 1.5)", "(0.6, 0.3)", "(5, 5)", "(1.5, 0.8)", "0.10"],
        ["Clutter", "(2, 1)", "(0.4, 0.3)", "(10, 8)", "(0.1, 0.1)", "0.15"],
    ]
)

# Slide 22: Classifier Comparison
add_comparison_slide(
    prs,
    "Rule-Based vs Naive Bayes Comparison",
    "Rule-Based Classifier",
    [
        "Fast computation (direct threshold checks)",
        "Highly interpretable (explicit rules)",
        "No training required",
        "Works well with domain expertise",
        "Binary matching (in range or not)",
        "Sensitive to threshold selection",
        "No probabilistic confidence",
        "Manual tuning for new scenarios"
    ],
    "Naive Bayes Classifier",
    [
        "Moderate computation (Gaussian PDF)",
        "Probabilistic output (soft decisions)",
        "Requires training data / parameters",
        "Data-driven parameter estimation",
        "Smooth probability gradients",
        "Handles feature overlap gracefully",
        "True confidence probabilities",
        "Adaptable through parameter updates"
    ]
)

# Slide 23: Full Pipeline Flow
add_flowchart_slide(prs, "Complete Classification Pipeline", [
    ("Detections", "Raw radar detections with (x, y, vr, RCS)"),
    ("DBSCAN", "Cluster detections into objects"),
    ("Tracking", "Associate clusters, maintain track IDs"),
    ("Features", "Extract cluster features for classification"),
    ("Classify", "Apply Rule-Based or Naive Bayes model")
])

# Slide 24: Feature History & Temporal Smoothing
add_content_slide(prs, "Feature History & Temporal Smoothing", [
    ("Feature History Storage", [
        "Each track maintains last 20 frames of features",
        "Features stored per update: points, extent, RCS, speed, etc.",
        "Enables temporal analysis and trend detection"
    ]),
    ("Time-Averaged Features", [
        "Classification uses averaged features for stability",
        "avg_feature = mean(feature_history)",
        "Reduces frame-to-frame variation",
        "More reliable classification for established tracks"
    ]),
    ("Classification History", [
        "Track stores last 10 classifications",
        "Enables voting or consensus approaches",
        "Useful for detecting classification changes"
    ]),
    ("Benefits", [
        "Reduces single-frame noise effects",
        "More consistent predictions over time",
        "Better handling of partial occlusions"
    ])
])

# Section 5: Summary
add_section_slide(prs, "Summary & Key Takeaways", "SECTION 5")

# Slide 26: Algorithm Summary
add_content_slide(prs, "Algorithm Summary", [
    ("Radar Simulation", [
        "Physics-based power calculation with radar equation",
        "Realistic point cloud generation on visible surfaces",
        "Configurable noise models for range, azimuth, velocity"
    ]),
    ("DBSCAN Clustering", [
        "Density-based grouping without prior cluster count",
        "Separates objects from clutter/noise naturally",
        "ε = 2.0m, minPts = 2 for radar applications"
    ]),
    ("Object Tracking", [
        "Greedy nearest-neighbor association",
        "Linear motion prediction for state estimation",
        "Track lifecycle with confirmation and deletion"
    ]),
    ("Classification", [
        "Rule-Based: fast, interpretable, threshold-based",
        "Naive Bayes: probabilistic, adaptive, data-driven",
        "Feature-based: points, extent, RCS, speed"
    ])
])

# Final slide
add_title_slide(prs, "Thank You", "Questions?")

# Save the presentation
output_path = "/workspace/radar_algorithms_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
