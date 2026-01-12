"""
Create a High-Level Presentation for ZOPPLER SYSTEMS
Black and White Professional Theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Black and White Color Scheme
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(40, 40, 40)
MEDIUM_GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(200, 200, 200)
OFF_WHITE = RGBColor(245, 245, 245)

# Company branding
COMPANY_NAME = "ZOPPLER SYSTEMS"


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with black background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Black background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLACK
    background.line.fill.background()
    
    # Company name at top
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.5))
    company_frame = company_box.text_frame
    company_para = company_frame.paragraphs[0]
    company_para.text = COMPANY_NAME
    company_para.font.size = Pt(18)
    company_para.font.bold = True
    company_para.font.color.rgb = LIGHT_GRAY
    company_para.alignment = PP_ALIGN.LEFT
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(3), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_frame = subtitle_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = LIGHT_GRAY
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, section_title, section_number=""):
    """Add a section divider slide with black background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Black background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLACK
    background.line.fill.background()
    
    # Company name at top
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.5))
    company_frame = company_box.text_frame
    company_para = company_frame.paragraphs[0]
    company_para.text = COMPANY_NAME
    company_para.font.size = Pt(14)
    company_para.font.bold = True
    company_para.font.color.rgb = MEDIUM_GRAY
    company_para.alignment = PP_ALIGN.LEFT
    
    # Section number
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = section_number
        num_para.font.size = Pt(24)
        num_para.font.color.rgb = MEDIUM_GRAY
        num_para.alignment = PP_ALIGN.CENTER
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = section_title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items, two_column=False):
    """Add a content slide with bullet points - white background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Black title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLACK
    title_bar.line.fill.background()
    
    # Company name in title bar
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(3), Inches(0.3))
    company_frame = company_box.text_frame
    company_para = company_frame.paragraphs[0]
    company_para.text = COMPANY_NAME
    company_para.font.size = Pt(10)
    company_para.font.bold = True
    company_para.font.color.rgb = MEDIUM_GRAY
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    if two_column and len(content_items) == 2:
        # Two columns
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        col1_title, col1_items = content_items[0]
        p = left_frame.paragraphs[0]
        p.text = col1_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLACK
        
        for item in col1_items:
            p = left_frame.add_paragraph()
            if isinstance(item, tuple):
                p.text = f"• {item[0]}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(8)
                for sub in item[1]:
                    sp = left_frame.add_paragraph()
                    sp.text = f"   – {sub}"
                    sp.font.size = Pt(14)
                    sp.font.color.rgb = MEDIUM_GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(8)
        
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        col2_title, col2_items = content_items[1]
        p = right_frame.paragraphs[0]
        p.text = col2_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLACK
        
        for item in col2_items:
            p = right_frame.add_paragraph()
            if isinstance(item, tuple):
                p.text = f"• {item[0]}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(8)
                for sub in item[1]:
                    sp = right_frame.add_paragraph()
                    sp.text = f"   – {sub}"
                    sp.font.size = Pt(14)
                    sp.font.color.rgb = MEDIUM_GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(8)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        first = True
        for item in content_items:
            if first:
                p = content_frame.paragraphs[0]
                first = False
            else:
                p = content_frame.add_paragraph()
            
            if isinstance(item, tuple):
                p.text = f"• {item[0]}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = BLACK
                p.space_before = Pt(12)
                for sub in item[1]:
                    sp = content_frame.add_paragraph()
                    sp.text = f"   – {sub}"
                    sp.font.size = Pt(16)
                    sp.font.color.rgb = MEDIUM_GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = BLACK
                p.space_before = Pt(12)
    
    return slide


def add_highlight_slide(prs, title, highlights):
    """Add a slide with highlighted key points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Black title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLACK
    title_bar.line.fill.background()
    
    # Company name
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(3), Inches(0.3))
    company_frame = company_box.text_frame
    company_para = company_frame.paragraphs[0]
    company_para.text = COMPANY_NAME
    company_para.font.size = Pt(10)
    company_para.font.bold = True
    company_para.font.color.rgb = MEDIUM_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Highlight boxes
    num_items = len(highlights)
    cols = min(3, num_items)
    rows = (num_items + cols - 1) // cols
    
    box_width = Inches(3.8)
    box_height = Inches(2.0)
    h_gap = Inches(0.5)
    v_gap = Inches(0.3)
    
    total_width = cols * box_width + (cols - 1) * h_gap
    start_x = (prs.slide_width - total_width) / 2
    start_y = Inches(1.6)
    
    for i, (highlight_title, highlight_text) in enumerate(highlights):
        row = i // cols
        col = i % cols
        x = start_x + col * (box_width + h_gap)
        y = start_y + row * (box_height + v_gap)
        
        # Box with border
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = OFF_WHITE
        box.line.color.rgb = BLACK
        box.line.width = Pt(2)
        
        # Highlight title
        h_title = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), box_width - Inches(0.4), Inches(0.5))
        htf = h_title.text_frame
        hp = htf.paragraphs[0]
        hp.text = highlight_title
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = BLACK
        
        # Highlight text
        h_text = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), box_width - Inches(0.4), box_height - Inches(0.9))
        htf2 = h_text.text_frame
        htf2.word_wrap = True
        hp2 = htf2.paragraphs[0]
        hp2.text = highlight_text
        hp2.font.size = Pt(14)
        hp2.font.color.rgb = MEDIUM_GRAY
    
    return slide


def add_flowchart_slide(prs, title, steps):
    """Add a slide with a horizontal flowchart in black and white."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Black title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLACK
    title_bar.line.fill.background()
    
    # Company name
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(3), Inches(0.3))
    company_frame = company_box.text_frame
    company_para = company_frame.paragraphs[0]
    company_para.text = COMPANY_NAME
    company_para.font.size = Pt(10)
    company_para.font.bold = True
    company_para.font.color.rgb = MEDIUM_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Calculate positions
    num_steps = len(steps)
    box_width = Inches(2.2)
    box_height = Inches(1.2)
    arrow_width = Inches(0.6)
    total_width = num_steps * box_width + (num_steps - 1) * arrow_width
    start_x = (prs.slide_width - total_width) / 2
    y_pos = Inches(2.5)
    
    for i, (step_title, step_desc) in enumerate(steps):
        x_pos = start_x + i * (box_width + arrow_width)
        
        # Box - alternating black and white
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_pos, y_pos, box_width, box_height
        )
        if i % 2 == 0:
            box.fill.solid()
            box.fill.fore_color.rgb = BLACK
            box.line.fill.background()
            text_color = WHITE
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = BLACK
            box.line.width = Pt(2)
            text_color = BLACK
        
        # Step number
        num_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.1), box_width, Inches(0.4))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = f"0{i + 1}"
        np.font.size = Pt(12)
        np.font.bold = True
        np.font.color.rgb = text_color
        np.alignment = PP_ALIGN.CENTER
        
        # Step title
        title_tb = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.4), box_width, Inches(0.7))
        ttf = title_tb.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = step_title
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = text_color
        tp.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < num_steps - 1:
            arrow_x = x_pos + box_width + Inches(0.05)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + Inches(0.4), Inches(0.5), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MEDIUM_GRAY
            arrow.line.fill.background()
        
        # Description below
        desc_box = slide.shapes.add_textbox(x_pos, y_pos + box_height + Inches(0.2), box_width, Inches(2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = step_desc
        dp.font.size = Pt(11)
        dp.font.color.rgb = MEDIUM_GRAY
        dp.alignment = PP_ALIGN.CENTER
    
    return slide


def add_closing_slide(prs, main_text, contact_info=""):
    """Add a closing/thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Black background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLACK
    background.line.fill.background()
    
    # Company name large
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
    company_frame = company_box.text_frame
    company_para = company_frame.paragraphs[0]
    company_para.text = COMPANY_NAME
    company_para.font.size = Pt(48)
    company_para.font.bold = True
    company_para.font.color.rgb = WHITE
    company_para.alignment = PP_ALIGN.CENTER
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.1), Inches(4.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()
    
    # Main text
    main_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    main_frame = main_box.text_frame
    main_para = main_frame.paragraphs[0]
    main_para.text = main_text
    main_para.font.size = Pt(28)
    main_para.font.color.rgb = LIGHT_GRAY
    main_para.alignment = PP_ALIGN.CENTER
    
    if contact_info:
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1))
        contact_frame = contact_box.text_frame
        contact_para = contact_frame.paragraphs[0]
        contact_para.text = contact_info
        contact_para.font.size = Pt(16)
        contact_para.font.color.rgb = MEDIUM_GRAY
        contact_para.alignment = PP_ALIGN.CENTER
    
    return slide


# =============================================================================
# CREATE THE PRESENTATION
# =============================================================================

# Slide 1: Title
add_title_slide(
    prs,
    "Traffic Radar Simulation",
    "Advanced Detection & Classification Platform"
)

# Slide 2: Overview
add_content_slide(prs, "Platform Overview", [
    ("Radar Detection Simulation", [
        "Physics-based signal modeling",
        "Realistic point cloud generation",
        "Configurable noise parameters"
    ]),
    ("Intelligent Clustering", [
        "DBSCAN density-based algorithm",
        "Automatic object separation",
        "Noise filtering"
    ]),
    ("Multi-Object Tracking", [
        "Real-time track management",
        "Motion prediction",
        "Data association"
    ]),
    ("Classification Engine", [
        "Rule-based classification",
        "Probabilistic (Naive Bayes) models",
        "Multi-class support"
    ])
])

# Section 1: Core Technology
add_section_slide(prs, "Core Technology", "01")

# Slide 4: Key Capabilities
add_highlight_slide(prs, "Key Capabilities", [
    ("Real-Time Processing", "High-performance algorithms optimized for real-time radar data processing and analysis"),
    ("Multi-Target Support", "Simultaneous tracking and classification of multiple traffic participants"),
    ("Flexible Output", "Support for UDP, Excel, and CSV output formats for system integration"),
    ("3D Visualization", "OpenGL-based visualization for scenario monitoring and validation"),
    ("Configurable Models", "Adjustable parameters for different radar configurations and scenarios"),
    ("Validated Algorithms", "Industry-standard algorithms with proven accuracy and reliability")
])

# Slide 5: System Architecture
add_flowchart_slide(prs, "Processing Pipeline", [
    ("Detection", "Raw radar data acquisition"),
    ("Clustering", "DBSCAN object grouping"),
    ("Tracking", "Multi-object track management"),
    ("Classification", "Object type identification"),
    ("Output", "Data export & visualization")
])

# Section 2: Supported Objects
add_section_slide(prs, "Traffic Participants", "02")

# Slide 7: Object Types
add_content_slide(prs, "Supported Object Classes", [
    ("Vehicles", [
        "Passenger cars",
        "Trucks and heavy vehicles",
        "Motorcycles and scooters"
    ]),
    ("Vulnerable Road Users", [
        "Pedestrians",
        "Cyclists",
        "E-scooters"
    ]),
    ("Environmental", [
        "Static clutter filtering",
        "Dynamic interference handling"
    ])
], two_column=False)

# Section 3: Applications
add_section_slide(prs, "Applications", "03")

# Slide 9: Use Cases
add_highlight_slide(prs, "Industry Applications", [
    ("ADAS Development", "Validation and testing of Advanced Driver Assistance Systems"),
    ("Autonomous Vehicles", "Sensor simulation for self-driving vehicle development"),
    ("Traffic Monitoring", "Infrastructure-based traffic flow analysis"),
    ("Safety Systems", "Collision warning and pedestrian detection validation"),
    ("Research & Academia", "Algorithm development and benchmarking"),
    ("Regulatory Compliance", "Testing against safety standards and regulations")
])

# Section 4: Benefits
add_section_slide(prs, "Benefits", "04")

# Slide 11: Key Benefits
add_content_slide(prs, "Why Choose ZOPPLER SYSTEMS", [
    ("Accuracy", [
        "Physics-based simulation models",
        "Validated against real-world data",
        "Comprehensive noise modeling"
    ]),
    ("Flexibility", [
        "Configurable radar parameters",
        "Multiple output formats",
        "Extensible architecture"
    ]),
    ("Performance", [
        "Real-time processing capability",
        "Efficient algorithms",
        "Scalable design"
    ]),
    ("Integration", [
        "Standard data interfaces",
        "API support",
        "Easy deployment"
    ])
])

# Final slide
add_closing_slide(
    prs,
    "Thank You",
    "www.zopplersystems.com | contact@zopplersystems.com"
)

# Save the presentation
output_path = "/workspace/zoppler_systems_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
